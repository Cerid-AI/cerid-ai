# Copyright (c) 2026 Justin Michaels. All rights reserved.
# SPDX-License-Identifier: FSL-1.1-ALv2

"""Office document parsers — DOCX, XLSX, and PPTX.

PPTX (C2.5) — minimal text-frame extraction via ``python-pptx``. One
``--- Slide N ---`` block per slide; shape text and notes are
concatenated. Legacy ``.ppt`` (CFB binary) is NOT supported by
python-pptx — uploads are rejected with a clear 422.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from fastapi import HTTPException

from app.parsers._zip_safety import assert_safe_zip
from app.parsers.registry import _MAX_TEXT_CHARS, logger, register_parser


@register_parser([".docx"])
def parse_docx(file_path: str) -> dict[str, Any]:
    assert_safe_zip(file_path)

    import docx

    try:
        doc = docx.Document(file_path)
    except HTTPException:
        raise
    except Exception as e:
        raise ValueError(
            f"Failed to read DOCX '{Path(file_path).name}': {e}. "
            f"File may be corrupted or not a valid .docx file."
        ) from e

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    table_texts = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            table_texts.append("\n".join(rows))

    parts = paragraphs
    if table_texts:
        parts.append("\n--- Tables ---")
        parts.extend(table_texts)

    text = "\n\n".join(parts)
    return {
        "text": text[:_MAX_TEXT_CHARS],
        "file_type": "docx",
        "page_count": None,
    }


@register_parser([".xlsx"])
def parse_xlsx(file_path: str) -> dict[str, Any]:
    """Parse XLSX with header auto-detection and Markdown table formatting."""
    assert_safe_zip(file_path)

    from openpyxl import load_workbook

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
    except HTTPException:
        raise
    except Exception as e:
        raise ValueError(
            f"Failed to read XLSX '{Path(file_path).name}': {e}. "
            f"File may be corrupted or not a valid .xlsx file."
        ) from e

    sheet_names = list(wb.sheetnames)
    sheets_text = []
    all_columns: list[str] = []
    total_rows = 0
    truncated = False

    for sheet_name in sheet_names:
        ws = wb[sheet_name]
        raw_rows = []
        for row in ws.iter_rows(values_only=True):
            cell_values = [str(c) if c is not None else "" for c in row]
            if any(cell_values):
                raw_rows.append(cell_values)

        if not raw_rows:
            continue

        total_rows += len(raw_rows)

        header_idx = 0
        for idx, row in enumerate(raw_rows[:5]):
            non_empty = sum(1 for c in row if c.strip())
            if len(row) > 0 and non_empty / len(row) > 0.5:
                header_idx = idx
                break

        header = raw_rows[header_idx]
        all_columns.extend([c.strip() for c in header if c.strip()])
        data_rows = raw_rows[header_idx + 1:]

        if len(data_rows) > 5000:
            logger.warning(
                f"XLSX '{Path(file_path).name}' sheet '{sheet_name}' has "
                f"{len(data_rows)} rows, truncating to 5000"
            )
            data_rows = data_rows[:5000]
            truncated = True

        md_lines = [f"--- Sheet: {sheet_name} ---"]
        md_lines.append("| " + " | ".join(header) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for row in data_rows:
            # Pad or trim row to match header column count
            padded = row + [""] * max(0, len(header) - len(row))
            md_lines.append("| " + " | ".join(padded[:len(header)]) + " |")
        sheets_text.append("\n".join(md_lines))

    wb.close()

    text = "\n\n".join(sheets_text)

    result: dict[str, Any] = {
        "text": text[:_MAX_TEXT_CHARS],
        "file_type": "xlsx",
        "page_count": len(sheet_names),
        "row_count": total_rows,
    }
    if all_columns:
        seen: set = set()
        unique_cols = []
        for c in all_columns:
            if c not in seen:
                seen.add(c)
                unique_cols.append(c)
        result["columns"] = json.dumps(unique_cols[:50])
    if truncated:
        result["truncated"] = True

    return result


@register_parser([".pptx"])
def parse_pptx(file_path: str) -> dict[str, Any]:
    """Parse PPTX — concatenate shape text + speaker notes per slide.

    Returns ``{text, file_type, page_count, slide_count}``.

    Matches the contract of :func:`parse_docx` — ``page_count`` is the
    slide count (peer parsers stuff their canonical "unit count" here)
    and ``slide_count`` is added as an explicit alias for retrieval-side
    consumers that want unambiguous slide semantics.
    """
    assert_safe_zip(file_path)

    from pptx import Presentation

    try:
        prs = Presentation(file_path)
    except HTTPException:
        raise
    except Exception as e:
        raise ValueError(
            f"Failed to read PPTX '{Path(file_path).name}': {e}. "
            f"File may be corrupted or not a valid .pptx file."
        ) from e

    sections: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"--- Slide {idx} ---"]

        for shape in slide.shapes:
            # ``has_text_frame`` short-circuits non-text shapes (pictures,
            # placeholders without text, group shapes' container nodes).
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if line:
                    parts.append(line)

        # Speaker notes — pptx exposes them via a separate slide.
        notes_slide = getattr(slide, "notes_slide", None) if slide.has_notes_slide else None
        if notes_slide is not None:
            notes_tf = getattr(notes_slide, "notes_text_frame", None)
            if notes_tf is not None:
                notes_text = (notes_tf.text or "").strip()
                if notes_text:
                    parts.append(f"[Notes] {notes_text}")

        # Skip slides that had no extractable text at all.
        if len(parts) > 1:
            sections.append("\n".join(parts))

    slide_count = len(prs.slides)
    text = "\n\n".join(sections)
    return {
        "text": text[:_MAX_TEXT_CHARS],
        "file_type": "pptx",
        "page_count": slide_count,
        "slide_count": slide_count,
    }


@register_parser([".ppt"])
def parse_ppt(file_path: str) -> dict[str, Any]:
    """Legacy ``.ppt`` (CFB binary) — unsupported; raise 422.

    python-pptx only handles the OOXML ``.pptx`` format. Surfacing this
    as a clear HTTP 422 lets the upload UI nudge the user to re-save as
    .pptx rather than crashing inside the parser.
    """
    raise HTTPException(
        status_code=422,
        detail=(
            f"Legacy .ppt format is not supported (file: '{Path(file_path).name}'). "
            "Please convert to .pptx in PowerPoint or Keynote and re-upload."
        ),
    )
